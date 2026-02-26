"""
VoleyPlay — Generador de presentación PPTX
Requiere: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Paleta de colores ───────────────────────────────────────────────
AZUL_OSCURO = RGBColor(0x0D, 0x2B, 0x55)   # #0D2B55
AZUL_MEDIO = RGBColor(0x1A, 0x6B, 0xA8)   # #1A6BA8
NARANJA = RGBColor(0xFF, 0x6B, 0x00)   # #FF6B00
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO = RGBColor(0xF2, 0xF4, 0xF7)
GRIS_TEXTO = RGBColor(0x44, 0x44, 0x55)
VERDE = RGBColor(0x2E, 0xCC, 0x71)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completamente en blanco


# ─── Helpers ─────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=BLANCO,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, color=GRIS_TEXTO, bullet="•"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def slide_bg(slide, color=AZUL_OSCURO):
    """Fondo de toda la diapositiva."""
    add_rect(slide, 0, 0, 13.33, 7.5, color)


def accent_bar(slide, color=NARANJA, height=0.08):
    """Barra de acento naranja en la parte superior."""
    add_rect(slide, 0, 0, 13.33, height, color)


def footer(slide, text="VoleyPlay — © 2026 Abdón Hernández Perera"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, AZUL_OSCURO)
    add_text(slide, text, 0.3, 7.12, 12, 0.35,
             font_size=10, color=RGBColor(0xAA, 0xBB, 0xCC),
             align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, AZUL_OSCURO)

# Gradiente lateral izquierdo (rectángulo decorativo)
add_rect(sl, 0, 0, 0.5, 7.5, NARANJA)
add_rect(sl, 0.5, 0, 0.12, 7.5, AZUL_MEDIO)

# Pelota de voleibol (círculos decorativos)
add_rect(sl, 10.5, 1.0, 2.5, 2.5, AZUL_MEDIO)
add_rect(sl, 10.8, 1.3, 1.9, 1.9, NARANJA)

# Título principal
add_text(sl, "🏐 VoleyPlay", 1.0, 1.8, 9, 1.5,
         font_size=60, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)

# Subtítulo
add_text(sl, "Aprende voleibol jugando",
         1.0, 3.2, 9, 0.8, font_size=28, color=RGBColor(0xAA, 0xCC, 0xFF),
         align=PP_ALIGN.LEFT)

# Descripción corta
add_text(sl,
         "Plataforma educativa para equipos juveniles de voleibol.\n"
         "Quiz · Desafío de campo · Simulador de rotaciones · Ranking global",
         1.0, 4.1, 10.5, 1.2, font_size=18,
         color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)

# Metadatos
add_text(sl, "Versión 1.0  ·  Febrero 2026",
         1.0, 6.3, 9, 0.5, font_size=14,
         color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problema / Contexto
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, GRIS_CLARO)
accent_bar(sl, AZUL_OSCURO, 0.12)
add_rect(sl, 0, 0.12, 0.08, 7.38, NARANJA)

add_text(sl, "El Problema", 0.4, 0.2, 12, 0.9,
         font_size=34, bold=True, color=AZUL_OSCURO)

# Columna izquierda
add_rect(sl, 0.4, 1.3, 5.8, 5.5, BLANCO)
add_text(sl, "Desafíos actuales", 0.7, 1.45, 5.2, 0.6,
         font_size=18, bold=True, color=AZUL_OSCURO)
add_bullet_box(sl, [
    "Los jugadores principiantes carecen de recursos\ninteractivos para aprender reglas y tácticas",
    "El entrenamiento tradicional es estático y poco\nmotivador para equipos juveniles",
    "No existe una herramienta que combine juego\neducativo + seguimiento de progreso + ranking",
    "Los entrenadores no pueden personalizar los\nejercizios de posicionamiento fácilmente",
], 0.7, 2.1, 5.3, 4.2, font_size=15, color=GRIS_TEXTO)

# Columna derecha
add_rect(sl, 6.9, 1.3, 6.0, 5.5, AZUL_OSCURO)
add_text(sl, "Nuestra Solución", 7.2, 1.45, 5.5, 0.6,
         font_size=18, bold=True, color=NARANJA)
add_bullet_box(sl, [
    "Quizzes interactivos con 28 preguntas en\n5 categorías temáticas",
    "Minijuegos de posicionamiento con\ndrag & drop en cancha virtual",
    "Simulador de rotaciones K1/K2 con\nvalidación en tiempo real",
    "Ranking global competitivo para motivar\nel aprendizaje continuo",
    "Panel de entrenador para personalizar\nlos ejercicios del equipo",
], 7.2, 2.1, 5.5, 4.2, font_size=15, color=RGBColor(0xCC, 0xDD, 0xEE))

footer(sl)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Funcionalidades principales
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, GRIS_CLARO)
accent_bar(sl, NARANJA, 0.12)
add_rect(sl, 0, 0.12, 0.08, 7.38, AZUL_OSCURO)

add_text(sl, "Funcionalidades Principales", 0.4, 0.2, 12, 0.9,
         font_size=34, bold=True, color=AZUL_OSCURO)

cards = [
    ("🧠", "Quiz de\nPreguntas",
     "28 preguntas · 5 categorías\nVelocidad y racha de bonos\nOpción múltiple con 30 s/pregunta",  AZUL_OSCURO),
    ("🏐", "Desafío de\nCampo",
     "Drag & Drop de jugadores\nValidación de posiciones\nConfiguración por entrenador",         AZUL_MEDIO),
    ("🔄", "Simulador de\nRotaciones",
     "Sistemas K1 y K2\nRondas de rotación interactivas\nFeedback inmediato",              NARANJA),
    ("🏆", "Ranking\nGlobal",     "Leaderboard en tiempo real\nPuntos acumulativos\nRacha de días activos",
     RGBColor(0x27, 0xAE, 0x60)),
]

card_w = 2.9
card_h = 4.8
start_x = 0.35
for i, (icon, title, body, color) in enumerate(cards):
    x = start_x + i * (card_w + 0.28)
    add_rect(sl, x, 1.3, card_w, card_h, color)
    add_text(sl, icon, x, 1.45, card_w, 0.9, font_size=36,
             align=PP_ALIGN.CENTER, color=BLANCO)
    add_text(sl, title, x, 2.3, card_w, 0.9, font_size=17, bold=True,
             align=PP_ALIGN.CENTER, color=BLANCO)
    add_rect(sl, x + 0.15, 3.25, card_w - 0.3, 0.04, BLANCO)
    add_bullet_box(sl, body.split("\n"), x + 0.15, 3.45,
                   card_w - 0.3, 2.5, font_size=13,
                   color=RGBColor(0xDD, 0xEE, 0xFF), bullet="›")

footer(sl)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Sistema de Puntuación
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, AZUL_OSCURO)
accent_bar(sl, NARANJA, 0.12)

add_text(sl, "Sistema de Puntuación", 0.5, 0.25, 12, 0.9,
         font_size=34, bold=True, color=BLANCO)

rows = [
    ("✅ Respuesta correcta",             "+100 pts", VERDE),
    ("⚡ Bono de velocidad  (< 10 s)",    "+50 pts",  NARANJA),
    ("🔥 Bono de racha  (3 correctas)",   "+50 pts",  RGBColor(0xE7, 0x4C, 0x3C)),
    ("💥 Bono de racha extendida (× 6)", "+100 pts", RGBColor(0x9B, 0x59, 0xB6)),
    ("❌ Respuesta incorrecta / tiempo",  "0 pts",    RGBColor(0x77, 0x77, 0x88)),
]

row_h = 0.88
start_y = 1.35
for i, (evento, pts, color) in enumerate(rows):
    y = start_y + i * (row_h + 0.1)
    add_rect(sl, 0.5, y, 0.08, row_h, color)
    add_rect(sl, 0.7, y, 8.8, row_h, RGBColor(0x1A, 0x2E, 0x4A))
    add_text(sl, evento, 1.0, y + 0.15, 7.5, 0.6,
             font_size=20, bold=True, color=BLANCO)
    add_rect(sl, 9.7, y, 2.9, row_h, color)
    add_text(sl, pts, 9.7, y + 0.15, 2.9, 0.6,
             font_size=22, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# Nota
add_text(sl,
         "Los puntos se acumulan permanentemente y determinan la posición en el ranking global.",
         0.5, 6.6, 12.3, 0.5, font_size=13,
         color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Arquitectura Técnica
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, GRIS_CLARO)
accent_bar(sl, AZUL_OSCURO, 0.12)
add_rect(sl, 0, 0.12, 0.08, 7.38, NARANJA)

add_text(sl, "Arquitectura Técnica", 0.4, 0.2, 12, 0.9,
         font_size=34, bold=True, color=AZUL_OSCURO)

layers = [
    ("Frontend",  "Angular 19 · Standalone Components\nAngular Material 19 · Facade Pattern\nLazy-loading · RxJS",
     AZUL_MEDIO,   0.4,  1.35),
    ("Backend",   "NestJS 10 · TypeORM 0.3\nJWT httpOnly Cookies · Passport\nHelmet · Rate Limiting · bcryptjs",
     AZUL_OSCURO,  4.65, 1.35),
    ("Base de Datos", "PostgreSQL 16\nMigraciones explícitas TypeORM\n4 entidades principales",
     RGBColor(0x2C, 0x3E, 0x50), 8.9, 1.35),
]

for title, body, color, lx, ly in layers:
    add_rect(sl, lx, ly, 4.0, 3.8, color)
    add_text(sl, title, lx + 0.15, ly + 0.15, 3.7, 0.65,
             font_size=19, bold=True, color=BLANCO)
    add_rect(sl, lx + 0.15, ly + 0.75, 3.7, 0.04, BLANCO)
    add_bullet_box(sl, body.split("\n"), lx + 0.2, ly + 0.95,
                   3.6, 2.8, font_size=15,
                   color=RGBColor(0xCC, 0xDD, 0xEE))

# Flechas (texto simulado)
add_text(sl, "→", 4.28, 2.7, 0.5, 0.5,
         font_size=30, bold=True, color=NARANJA, align=PP_ALIGN.CENTER)
add_text(sl, "→", 8.55, 2.7, 0.5, 0.5,
         font_size=30, bold=True, color=NARANJA, align=PP_ALIGN.CENTER)

# Infraestructura
add_rect(sl, 0.4, 5.4, 12.5, 1.35, AZUL_OSCURO)
add_text(sl, "🐳  Infraestructura Docker",
         0.65, 5.5, 4, 0.5, font_size=16, bold=True, color=NARANJA)
add_bullet_box(sl, [
    "docker-compose (dev) · docker-compose.prod.yml (producción)",
    "Contenedores: voley_db · voley_backend · voley_frontend  |  Usuario no-root appuser",
], 0.65, 6.0, 12, 0.7, font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

footer(sl)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Seguridad
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, AZUL_OSCURO)
accent_bar(sl, NARANJA, 0.12)

add_text(sl, "Seguridad", 0.5, 0.25, 12, 0.9,
         font_size=34, bold=True, color=BLANCO)

items = [
    ("🔐", "JWT httpOnly",    "Token almacenado exclusivamente en cookie httpOnly Secure.\nSin acceso desde JavaScript del navegador."),
    ("🛡️", "Helmet",          "Headers HTTP de seguridad: CSP, HSTS, X-Frame-Options, etc."),
    ("⚡", "Rate Limiting",   "120 req/min global · 8 req/min en endpoints de autenticación."),
    ("🔑", "bcryptjs",        "12 rondas de hashing para contraseñas de usuario."),
    ("🌐", "CORS estricto",   "Solo orígenes explícitamente configurados en CORS_ORIGINS."),
    ("🚫", "CSRF Protection",
     "Middleware de verificación de origen en todas las peticiones."),
]

col_w = 5.8
for i, (icon, title, body) in enumerate(items):
    col = i % 2
    row = i // 2
    x = 0.4 + col * (col_w + 0.75)
    y = 1.45 + row * 1.85
    add_rect(sl, x, y, col_w, 1.65, RGBColor(0x0A, 0x1E, 0x3C))
    add_text(sl, icon, x + 0.15, y + 0.1, 0.8,
             0.8, font_size=26, color=NARANJA)
    add_text(sl, title, x + 1.0, y + 0.1, col_w - 1.15, 0.5,
             font_size=17, bold=True, color=BLANCO)
    add_text(sl, body, x + 1.0, y + 0.6, col_w - 1.15, 0.9,
             font_size=13, color=RGBColor(0xAA, 0xBB, 0xCC))

add_text(sl,
         "Producción: validación de env.config.ts rechaza JWT_SECRET débil y CORS no seguros.",
         0.5, 7.05, 12.3, 0.4, font_size=12,
         color=RGBColor(0x66, 0x77, 0x88), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Panel de Entrenador
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, GRIS_CLARO)
accent_bar(sl, NARANJA, 0.12)
add_rect(sl, 0, 0.12, 0.08, 7.38, AZUL_OSCURO)

add_text(sl, "Panel de Entrenador", 0.4, 0.2, 12, 0.9,
         font_size=34, bold=True, color=AZUL_OSCURO)

# Columna izquierda
add_rect(sl, 0.4, 1.3, 5.9, 5.5, AZUL_OSCURO)
add_text(sl, "👨‍💼  Funciones exclusivas", 0.7, 1.45, 5.3, 0.7,
         font_size=18, bold=True, color=NARANJA)
add_bullet_box(sl, [
    "Acceso restringido: solo usuarios con rol coach",
    "Editor visual de zonas de formación",
    "Selección de modo: Field Challenge o Rotation",
    "Configuración por sistema: K1 / K2",
    "Drag & resize de rectángulos de zona",
    "Estadísticas de rendimiento del equipo",
    "Guardado persistente en la base de datos",
], 0.7, 2.2, 5.3, 4.3, font_size=15, color=RGBColor(0xCC, 0xDD, 0xEE))

# Columna derecha — flujo
add_rect(sl, 6.6, 1.3, 6.4, 5.5, BLANCO)
add_text(sl, "Flujo de configuración", 6.9, 1.45, 5.8, 0.6,
         font_size=17, bold=True, color=AZUL_OSCURO)

steps = [
    ("1", "Seleccionar modo de juego",    AZUL_MEDIO),
    ("2", "Elegir sistema K1 o K2",       AZUL_MEDIO),
    ("3", "Ajustar zonas en editor visual", NARANJA),
    ("4", "Guardar configuración",         VERDE),
    ("5", "¡Jugadores ya pueden practicar!", AZUL_OSCURO),
]
for i, (num, label, color) in enumerate(steps):
    y = 2.2 + i * 0.9
    add_rect(sl, 6.9, y, 0.5, 0.6, color)
    add_text(sl, num, 6.9, y, 0.5, 0.6, font_size=17, bold=True,
             color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(sl, label, 7.55, y + 0.05, 5.2, 0.55,
             font_size=15, color=GRIS_TEXTO)
    if i < len(steps) - 1:
        add_text(sl, "↓", 7.1, y + 0.62, 0.5, 0.3,
                 font_size=14, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)

footer(sl)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Stack Tecnológico
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, AZUL_OSCURO)
accent_bar(sl, NARANJA, 0.12)

add_text(sl, "Stack Tecnológico", 0.5, 0.25, 12, 0.9,
         font_size=34, bold=True, color=BLANCO)

techs = [
    # (nombre, versión, categoría, color)
    ("Angular",       "19",    "Frontend",       AZUL_MEDIO),
    ("Ang. Material", "19",    "UI Library",      RGBColor(0x0F, 0x9D, 0x58)),
    ("TypeScript",    "5.x",   "Lenguaje",        RGBColor(0x30, 0x78, 0xC5)),
    ("NestJS",        "10",    "Backend",         RGBColor(0xE0, 0x22, 0x4C)),
    ("TypeORM",       "0.3",   "ORM",             RGBColor(0x66, 0x44, 0xCC)),
    ("PostgreSQL",    "16",    "Base de Datos",   RGBColor(0x33, 0x6E, 0x9E)),
    ("Docker",        "24+",   "Infraestructura", RGBColor(0x00, 0x91, 0xE2)),
    ("JWT",           "Cookies", "Seguridad",       RGBColor(0xD3, 0x54, 0x00)),
]

card_w = 2.9
card_h = 1.8
cols = 4
for i, (name, version, cat, color) in enumerate(techs):
    col = i % cols
    row = i // cols
    x = 0.35 + col * (card_w + 0.35)
    y = 1.5 + row * (card_h + 0.3)
    add_rect(sl, x, y, card_w, card_h, color)
    add_text(sl, name, x + 0.1, y + 0.15, card_w - 0.2, 0.65,
             font_size=20, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(sl, f"v{version}", x + 0.1, y + 0.8, card_w - 0.2, 0.45,
             font_size=16, color=BLANCO, align=PP_ALIGN.CENTER)
    add_rect(sl, x + 0.1, y + 1.3, card_w - 0.2, 0.04, BLANCO)
    add_text(sl, cat, x + 0.1, y + 1.38, card_w - 0.2, 0.35,
             font_size=12, color=RGBColor(0xDD, 0xEE, 0xFF),
             align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Instalación rápida
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, GRIS_CLARO)
accent_bar(sl, AZUL_OSCURO, 0.12)
add_rect(sl, 0, 0.12, 0.08, 7.38, NARANJA)

add_text(sl, "Instalación Rápida", 0.4, 0.2, 12, 0.9,
         font_size=34, bold=True, color=AZUL_OSCURO)

# Código
code_lines = [
    "# 1. Clonar el repositorio",
    "git clone <repo-url> voley-app && cd voley-app",
    "",
    "# 2. Levantar todos los servicios con Docker",
    "docker-compose up -d",
    "",
    "# 3. Cargar preguntas iniciales (solo primera vez)",
    "docker exec voley_backend sh -c 'npm run seed'",
    "",
    "# 4. Abrir la aplicación",
    "open http://localhost:4200",
]

add_rect(sl, 0.4, 1.3, 12.5, 5.2, RGBColor(0x1E, 0x1E, 0x2E))
y_cursor = 1.5
for line in code_lines:
    color = RGBColor(0x88, 0xCC, 0x88) if line.startswith("#") else (
        RGBColor(0xFF, 0xCC, 0x66) if line.startswith("docker") or line.startswith("git") or line.startswith("open") else
        RGBColor(0xDD, 0xDD, 0xEE))
    if line:
        add_text(sl, line, 0.75, y_cursor, 12, 0.38,
                 font_size=16, color=color, align=PP_ALIGN.LEFT)
    y_cursor += 0.42

footer(sl)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Roadmap / Próximas características
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, AZUL_OSCURO)
accent_bar(sl, NARANJA, 0.12)

add_text(sl, "Roadmap", 0.5, 0.25, 12, 0.9,
         font_size=34, bold=True, color=BLANCO)

phases = [
    ("✅ Fase 1 — Completada",
     ["Base de datos + migraciones", "Sistema de autenticación JWT",
         "Quiz de preguntas (28 items)", "Ranking global", "Panel de entrenador básico"],
     VERDE, 0.4),
    ("✅ Fase 2 — Completada",
     ["Editor de formaciones visual", "Desafío de campo configurable",
         "Simulador K1/K2", "Zonas de posición en BD", "Tests unitarios (45 tests)"],
     AZUL_MEDIO, 4.6),
    ("🔜 Fase 3 — Próximamente",
     ["App móvil nativa (Ionic)", "Notificaciones push de racha", "Modo multijugador en tiempo real",
      "Video-tutoriales integrados", "Exportación de estadísticas PDF"],
     NARANJA, 8.8),
]

for title, items, color, x in phases:
    add_rect(sl, x, 1.4, 4.0, 5.5, RGBColor(0x0A, 0x1E, 0x3C))
    add_rect(sl, x, 1.4, 4.0, 0.65, color)
    add_text(sl, title, x + 0.1, 1.45, 3.8, 0.55,
             font_size=14, bold=True, color=BLANCO)
    add_bullet_box(sl, items, x + 0.15, 2.15, 3.7, 4.5,
                   font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Cierre
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_bg(sl, AZUL_OSCURO)

# Franja naranja inferior
add_rect(sl, 0, 5.8, 13.33, 1.7, NARANJA)

# Decorativo
add_rect(sl, 0, 0, 13.33, 0.08, NARANJA)
add_rect(sl, 6.4, 0, 0.08, 7.5, RGBColor(0x1A, 0x3A, 0x6B))

add_text(sl, "🏐", 0.5, 0.8, 5.8, 2.5,
         font_size=100, align=PP_ALIGN.CENTER, color=BLANCO)

add_text(sl, "VoleyPlay", 6.7, 1.2, 6.3, 1.4,
         font_size=52, bold=True, color=BLANCO)

add_text(sl, "Aprende voleibol jugando.", 6.7, 2.7, 6.3, 0.8,
         font_size=24, color=RGBColor(0xAA, 0xCC, 0xFF), italic=True)

add_text(sl, "📦  Repositorio:  github.com/tu-org/voley-app",
         6.7, 3.7, 6.3, 0.6, font_size=15, color=RGBColor(0x88, 0xAA, 0xCC))

add_text(sl, "¡Gracias!", 0.5, 5.9, 13.33 - 1, 1.0,
         font_size=40, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_text(sl, "Versión 1.0  ·  Febrero 2026  ·  Abdón Hernández Perera",
         0.5, 6.8, 13.33 - 1, 0.5, font_size=14,
         color=AZUL_OSCURO, align=PP_ALIGN.CENTER)


# ─── Guardar ──────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__),
                        "VoleyPlay_presentacion.pptx")
prs.save(out_path)
print(f"OK  Presentacion guardada en: {out_path}")
